#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
PowerPoint Translation Automation Script

This script automates the translation of PowerPoint (PPTX) files while 
preserving formatting and handling right-to-left (RTL) and left-to-right (LTR) 
language differences.

Usage:
    python pptx_translator.py input.pptx output.pptx en es
"""

import os
import asyncio
import argparse
import time
import logging
from typing import List, Dict, Tuple, Set, Any, Optional
from pathlib import Path
from dataclasses import dataclass
import pptx
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, MSO_TEXT_ORIENTATION
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.text.text import TextFrame, _Paragraph, _Run
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from concurrent.futures import ThreadPoolExecutor

# Import the translator module
from translator import AsyncTranslator, setup_logging, TranslationResult

# RTL language codes
RTL_LANGUAGES = {
    'ar', 'arc', 'dv', 'fa', 'ha', 'he', 'khw', 'ks', 'ku', 'ps', 'ur', 'yi'
}

# Default text directions by language ID
RTL_LANGUAGE_IDS = {
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_ALGERIA,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_BAHRAIN,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_EGYPT,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_IRAQ,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_JORDAN,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_KUWAIT,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_LEBANON,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_LIBYA,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_MOROCCO,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_OMAN,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_QATAR,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_SAUDI_ARABIA,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_SYRIA,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_TUNISIA,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_UAE,
    pptx.enum.lang.MSO_LANGUAGE_ID.ARABIC_YEMEN,
    pptx.enum.lang.MSO_LANGUAGE_ID.HEBREW,
    pptx.enum.lang.MSO_LANGUAGE_ID.PERSIAN,
    pptx.enum.lang.MSO_LANGUAGE_ID.URDU
}

# Language code to MSO_LANGUAGE_ID mapping
LANGUAGE_CODE_TO_MSO_ID = {
    'en': MSO_LANGUAGE_ID.ENGLISH_US,
    'es': MSO_LANGUAGE_ID.SPANISH,
    'fr': MSO_LANGUAGE_ID.FRENCH,
    'de': MSO_LANGUAGE_ID.GERMAN,
    'it': MSO_LANGUAGE_ID.ITALIAN,
    'pt': MSO_LANGUAGE_ID.PORTUGUESE_BRAZIL,
    'ru': MSO_LANGUAGE_ID.RUSSIAN,
    'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ko': MSO_LANGUAGE_ID.KOREAN,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'ar': MSO_LANGUAGE_ID.ARABIC,
    'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI,
    'fa': MSO_LANGUAGE_ID.PERSIAN,
    'ur': MSO_LANGUAGE_ID.URDU,
    # Add more language mappings as needed
}

@dataclass
class TextElement:
    """Class to store text elements with their metadata for translation"""
    id: str  # Unique identifier for tracking
    text: str  # Original text
    shape_id: Optional[int] = None  # ID of parent shape
    slide_num: int = 0  # Slide number (0-indexed)
    element_type: str = ""  # Type of element (e.g., 'paragraph', 'run', 'title')
    placeholders: List[Tuple[str, object]] = None  # Store references to original objects
    translated_text: Optional[str] = None  # Will store the translated text


class PPTXTranslator:
    """
    PowerPoint Translation Automation Class
    Handles extraction, translation, and reinsertion of text while
    preserving formatting and handling RTL/LTR language directions.
    """
    
    def __init__(self, logger=None, concurrent_slides=5, max_translation_batch=50):
        """
        Initialize the PowerPoint translator.
        
        Args:
            logger: Custom logger (if None, a default one will be created)
            concurrent_slides: Number of slides to process concurrently
            max_translation_batch: Maximum number of texts in one translation batch
        """
        self.logger = logger or setup_logging('pptx_translator.log', logging.INFO)
        self.translator = AsyncTranslator(logger=self.logger)
        self.concurrent_slides = concurrent_slides
        self.max_translation_batch = max_translation_batch
        self.text_elements = []  # Will store all extracted text elements
        self.element_count = 0  # Counter for generating unique IDs
        
    def _is_rtl_language(self, lang_code: str) -> bool:
        """Check if a language is RTL based on its code"""
        return lang_code.lower().split('-')[0] in RTL_LANGUAGES
    
    def _get_mso_language_id(self, lang_code: str) -> int:
        """Get the MSO_LANGUAGE_ID for a language code"""
        base_lang = lang_code.lower().split('-')[0]
        return LANGUAGE_CODE_TO_MSO_ID.get(base_lang, MSO_LANGUAGE_ID.ENGLISH_US)
    
    def extract_text_from_paragraph(self, paragraph: _Paragraph, slide_num: int, shape_id: int, element_type: str) -> List[TextElement]:
        """Extract text from a paragraph, potentially broken into runs"""
        elements = []
        
        # If the paragraph has a single run or consistent formatting, extract the whole paragraph
        if len(paragraph.runs) <= 1 or all(run.font.name == paragraph.runs[0].font.name for run in paragraph.runs):
            element_id = f"s{slide_num}_p{shape_id}_{self.element_count}"
            self.element_count += 1
            
            text = paragraph.text.strip()
            if text:
                elements.append(TextElement(
                    id=element_id,
                    text=text,
                    slide_num=slide_num,
                    shape_id=shape_id,
                    element_type=element_type,
                    placeholders=[(element_type, paragraph)]
                ))
        else:
            # Extract individual runs if they have different formatting
            for i, run in enumerate(paragraph.runs):
                text = run.text.strip()
                if text:
                    element_id = f"s{slide_num}_p{shape_id}_r{i}_{self.element_count}"
                    self.element_count += 1
                    elements.append(TextElement(
                        id=element_id,
                        text=text,
                        slide_num=slide_num,
                        shape_id=shape_id,
                        element_type="run",
                        placeholders=[("run", run)]
                    ))
        
        return elements
    
    def extract_text_from_shape(self, shape: BaseShape, slide_num: int) -> List[TextElement]:
        """Extract text from a PowerPoint shape"""
        elements = []
        
        # Skip shapes without text
        if not hasattr(shape, 'text_frame'):
            return elements
        
        try:
            # Get shape ID or generate one
            shape_id = getattr(shape, 'shape_id', hash(shape) % 10000)
            
            # Handle various types of shapes with text
            if hasattr(shape, 'text') and shape.text.strip():
                text_frame = getattr(shape, 'text_frame', None)
                
                # Extract from text frame with paragraphs
                if text_frame and hasattr(text_frame, 'paragraphs'):
                    for i, paragraph in enumerate(text_frame.paragraphs):
                        elements.extend(self.extract_text_from_paragraph(
                            paragraph, slide_num, shape_id, 
                            "title" if i == 0 and hasattr(shape, 'is_title') and shape.is_title else "paragraph"
                        ))
                # Extract from shape without proper text frame structure
                else:
                    element_id = f"s{slide_num}_sh{shape_id}_{self.element_count}"
                    self.element_count += 1
                    elements.append(TextElement(
                        id=element_id,
                        text=shape.text.strip(),
                        slide_num=slide_num,
                        shape_id=shape_id,
                        element_type="shape_text",
                        placeholders=[("shape", shape)]
                    ))
                
            # Handle tables
            if hasattr(shape, 'table'):
                table = shape.table
                for r_idx, row in enumerate(table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        for i, paragraph in enumerate(cell.text_frame.paragraphs):
                            element_id = f"s{slide_num}_t{shape_id}_r{r_idx}_c{c_idx}_p{i}_{self.element_count}"
                            self.element_count += 1
                            
                            text = paragraph.text.strip()
                            if text:
                                elements.append(TextElement(
                                    id=element_id,
                                    text=text,
                                    slide_num=slide_num,
                                    shape_id=shape_id,
                                    element_type="table_cell",
                                    placeholders=[("paragraph", paragraph)]
                                ))
        
        except Exception as e:
            self.logger.error(f"Error extracting text from shape on slide {slide_num}: {str(e)}")
        
        return elements
    
    def extract_text_from_slide(self, slide: Slide, slide_idx: int) -> List[TextElement]:
        """Extract all text elements from a slide"""
        elements = []
        
        # Process all shapes in the slide
        for shape in slide.shapes:
            elements.extend(self.extract_text_from_shape(shape, slide_idx))
            
            # Handle group shapes (shapes within shapes)
            if hasattr(shape, 'shapes'):
                for subshape in shape.shapes:
                    elements.extend(self.extract_text_from_shape(subshape, slide_idx))
        
        return elements
    
    async def extract_all_text(self, presentation: Presentation) -> List[TextElement]:
        """Extract all text from a PowerPoint presentation"""
        self.logger.info(f"Extracting text from {len(presentation.slides)} slides")
        self.text_elements = []
        self.element_count = 0
        
        # Process slides concurrently using a thread pool
        with ThreadPoolExecutor(max_workers=self.concurrent_slides) as executor:
            # Create tasks for each slide
            futures = []
            for slide_idx, slide in enumerate(presentation.slides):
                future = executor.submit(self.extract_text_from_slide, slide, slide_idx)
                futures.append(future)
            
            # Collect results as they complete
            for future in futures:
                elements = future.result()
                self.text_elements.extend(elements)
        
        self.logger.info(f"Extracted {len(self.text_elements)} text elements from presentation")
        return self.text_elements
    
    async def translate_text_elements(self, 
                                     text_elements: List[TextElement],
                                     source_lang: str,
                                     target_lang: str) -> List[TextElement]:
        """Translate all extracted text elements"""
        self.logger.info(f"Translating {len(text_elements)} text elements from {source_lang} to {target_lang}")
        
        # Prepare batches
        texts = [element.text for element in text_elements]
        
        # Translate in batches
        results = await self.translator.batch_translate_large(
            texts=texts,
            src_lang=source_lang,
            dest_lang=target_lang,
            batch_size=self.max_translation_batch
        )
        
        # Update text elements with translations
        for element, result in zip(text_elements, results):
            if result.success and result.translated_text:
                element.translated_text = result.translated_text
            else:
                self.logger.warning(f"Translation failed for element: {element.id}")
                # Keep original text as fallback
                element.translated_text = element.text
        
        return text_elements
    
    def apply_rtl_settings(self, element_type: str, obj: Any, is_rtl: bool):
        """Apply RTL/LTR text direction settings to PowerPoint objects"""
        try:
            if is_rtl:
                # Apply RTL settings based on element type
                if element_type in ("paragraph", "title", "table_cell") and hasattr(obj, 'alignment'):
                    # Right align for RTL languages
                    obj.alignment = pptx.enum.text.PP_ALIGN.RIGHT
                
                # Set paragraph direction for text frames
                if hasattr(obj, 'text_frame'):
                    text_frame = obj.text_frame
                    # Set RTL direction for text frame if possible
                    if hasattr(text_frame, '_element'):
                        # Add BiDi property to XML
                        p_pr = text_frame._element.get_or_add_pPr()
                        p_pr.set('rtl', '1')
                
                # For runs, try to set the RTL attribute
                if element_type == "run" and hasattr(obj, '_r'):
                    # Add RTL attribute to run properties
                    r_pr = obj._r.get_or_add_rPr()
                    r_pr.set('rtl', '1')
            else:
                # Apply LTR settings based on element type
                if element_type in ("paragraph", "title", "table_cell") and hasattr(obj, 'alignment'):
                    # Left align for LTR languages
                    obj.alignment = pptx.enum.text.PP_ALIGN.LEFT
                
                # Clear RTL setting if previously set
                if hasattr(obj, 'text_frame') and hasattr(obj.text_frame, '_element'):
                    p_pr = obj.text_frame._element.get_or_add_pPr()
                    if 'rtl' in p_pr.attrib:
                        del p_pr.attrib['rtl']
                
                if element_type == "run" and hasattr(obj, '_r'):
                    r_pr = obj._r.get_or_add_rPr()
                    if 'rtl' in r_pr.attrib:
                        del r_pr.attrib['rtl']
        except Exception as e:
            self.logger.error(f"Error applying RTL/LTR settings: {str(e)}")
    
    def update_language_settings(self, obj: Any, lang_id: int):
        """Update language settings of PowerPoint objects"""
        try:
            # For paragraphs and runs with language property
            if hasattr(obj, 'font') and hasattr(obj.font, 'language_id'):
                obj.font.language_id = lang_id
            
            # For shapes with text frames
            elif hasattr(obj, 'text_frame'):
                for paragraph in obj.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, 'font') and hasattr(run.font, 'language_id'):
                            run.font.language_id = lang_id
        except Exception as e:
            self.logger.error(f"Error updating language settings: {str(e)}")
    
    def apply_translations(self, 
                          text_elements: List[TextElement], 
                          target_lang: str):
        """Apply translated text back to the PowerPoint objects"""
        self.logger.info(f"Applying {len(text_elements)} translations to presentation")
        
        # Get target language properties
        is_rtl = self._is_rtl_language(target_lang)
        lang_id = self._get_mso_language_id(target_lang)
        
        # Group elements by slide for efficiency
        slides_map = {}
        for element in text_elements:
            if element.slide_num not in slides_map:
                slides_map[element.slide_num] = []
            slides_map[element.slide_num].append(element)
        
        # Process each slide
        for slide_num, elements in slides_map.items():
            self.logger.info(f"Applying translations to slide {slide_num} ({len(elements)} elements)")
            
            for element in elements:
                if not element.translated_text:
                    continue
                
                try:
                    # Apply translation to each placeholder
                    for ph_type, obj in element.placeholders:
                        if ph_type in ("paragraph", "title", "table_cell"):
                            # For paragraphs and table cells
                            if hasattr(obj, 'text'):
                                obj.text = element.translated_text
                                
                                # Apply RTL/LTR settings
                                self.apply_rtl_settings(ph_type, obj, is_rtl)
                                
                                # Update language settings
                                self.update_language_settings(obj, lang_id)
                                
                        elif ph_type == "run":
                            # For text runs
                            if hasattr(obj, 'text'):
                                obj.text = element.translated_text
                                
                                # Apply RTL/LTR settings
                                self.apply_rtl_settings(ph_type, obj, is_rtl)
                                
                                # Update language settings
                                if hasattr(obj, 'font') and hasattr(obj.font, 'language_id'):
                                    obj.font.language_id = lang_id
                                
                        elif ph_type == "shape":
                            # For shapes with direct text property
                            if hasattr(obj, 'text'):
                                obj.text = element.translated_text
                                
                                # Apply RTL/LTR settings if possible
                                if hasattr(obj, 'text_frame'):
                                    self.apply_rtl_settings("shape", obj, is_rtl)
                                    
                                    # Update language settings
                                    self.update_language_settings(obj, lang_id)
                
                except Exception as e:
                    self.logger.error(f"Error applying translation for element {element.id}: {str(e)}")
        
        self.logger.info("All translations applied successfully")
    
    def adjust_text_boxes_for_translations(self, presentation: Presentation):
        """Adjust text boxes to accommodate translations which might be longer or shorter"""
        self.logger.info("Adjusting text boxes for translated content")
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    try:
                        # Set autofit to accommodate text
                        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        
                        # Try to fit text by adjusting margins if needed
                        if len(shape.text) > 100:  # Only for longer text
                            shape.text_frame.margin_bottom = 0
                            shape.text_frame.margin_top = 0
                            shape.text_frame.margin_left = 0
                            shape.text_frame.margin_right = 0
                            
                    except Exception as e:
                        self.logger.warning(f"Error adjusting text box: {str(e)}")
    
    async def translate_presentation(self, 
                                    input_file: str, 
                                    output_file: str, 
                                    source_lang: str,
                                    target_lang: str) -> Dict[str, Any]:
        """
        Main method to translate a PowerPoint presentation
        
        Args:
            input_file: Path to input PPTX file
            output_file: Path to output PPTX file
            source_lang: Source language code
            target_lang: Target language code
            
        Returns:
            Dictionary with translation statistics
        """
        start_time = time.time()
        self.logger.info(f"Starting translation of '{input_file}' from {source_lang} to {target_lang}")
        
        try:
            # Load the presentation
            presentation = Presentation(input_file)
            
            # Extract text
            text_elements = await self.extract_all_text(presentation)
            
            if not text_elements:
                self.logger.warning("No text found to translate in the presentation")
                presentation.save(output_file)
                return {
                    "status": "success",
                    "elements_found": 0,
                    "elements_translated": 0,
                    "time_taken": time.time() - start_time
                }
            
            # Translate text
            translated_elements = await self.translate_text_elements(
                text_elements, 
                source_lang, 
                target_lang
            )
            
            # Apply translations back to the presentation
            self.apply_translations(translated_elements, target_lang)
            
            # Adjust text boxes if needed
            self.adjust_text_boxes_for_translations(presentation)
            
            # Save the translated presentation
            self.logger.info(f"Saving translated presentation to '{output_file}'")
            presentation.save(output_file)
            
            # Calculate statistics
            success_count = sum(1 for e in translated_elements if e.translated_text != e.text)
            
            stats = {
                "status": "success",
                "elements_found": len(text_elements),
                "elements_translated": success_count,
                "time_taken": time.time() - start_time
            }
            
            self.logger.info(f"Translation completed in {stats['time_taken']:.2f} seconds")
            self.logger.info(f"Translated {stats['elements_translated']} out of {stats['elements_found']} elements")
            
            return stats
            
        except Exception as e:
            self.logger.error(f"Error translating presentation: {str(e)}")
            return {
                "status": "error",
                "error": str(e),
                "time_taken": time.time() - start_time
            }

async def main():
    """Main entry point for the script"""
    # Set up argument parser
    parser = argparse.ArgumentParser(description='Translate PowerPoint presentations while preserving formatting.')
    parser.add_argument('input_file', help='Path to the input PPTX file')
    parser.add_argument('output_file', help='Path to save the translated PPTX file')
    parser.add_argument('source_lang', help='Source language code (e.g., "en", "auto")')
    parser.add_argument('target_lang', help='Target language code (e.g., "es", "ar")')
    parser.add_argument('--concurrent', type=int, default=5, help='Number of slides to process concurrently')
    parser.add_argument('--batch-size', type=int, default=50, help='Maximum batch size for translation')
    parser.add_argument('--log-level', choices=['DEBUG', 'INFO', 'WARNING', 'ERROR'], default='INFO', 
                        help='Set the logging level')
    
    args = parser.parse_args()
    
    # Set up logging
    log_level = getattr(logging, args.log_level)
    logger = setup_logging('pptx_translator.log', log_level)
    
    # Check if input file exists
    if not os.path.isfile(args.input_file):
        logger.error(f"Input file '{args.input_file}' not found")
        return
    
    # Create translator and translate
    translator = PPTXTranslator(
        logger=logger,
        concurrent_slides=args.concurrent,
        max_translation_batch=args.batch_size
    )
    
    logger.info(f"Starting translation: {args.input_file} ({args.source_lang}) -> {args.output_file} ({args.target_lang})")
    
    try:
        stats = await translator.translate_presentation(
            args.input_file,
            args.output_file,
            args.source_lang,
            args.target_lang
        )
        
        if stats["status"] == "success":
            logger.info("=" * 50)
            logger.info("Translation completed successfully")
            logger.info(f"Elements found: {stats['elements_found']}")
            logger.info(f"Elements translated: {stats['elements_translated']}")
            logger.info(f"Time taken: {stats['time_taken']:.2f} seconds")
            logger.info("=" * 50)
        else:
            logger.error(f"Translation failed: {stats.get('error', 'Unknown error')}")
    
    except Exception as e:
        logger.error(f"Error in translation process: {str(e)}")

if __name__ == "__main__":
    # Set up asyncio policies for Windows if needed
    if hasattr(asyncio, 'WindowsSelectorEventLoopPolicy') and hasattr(asyncio, 'set_event_loop_policy'):
        asyncio.set_event_loop_policy(asyncio.WindowsSelectorEventLoopPolicy())
    
    asyncio.run(main())
